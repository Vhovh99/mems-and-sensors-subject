"""Design system for the MEMS & Sensors lecture decks.

One place for palette, type scale, and slide furniture so both decks read as one course.
Diagrams are built from native PowerPoint shapes and stay editable in PowerPoint/Impress.

Semantic colour (used identically in every deck):
    TEAL   the true signal path, correct answers, the thing that works
    AMBER  where error enters
    RED    the term that kills the design / failure
    GRAY   context, secondary, de-emphasised
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# ---------------------------------------------------------------- palette
INK    = RGBColor(0x1A, 0x1F, 0x24)   # primary text
GROUND = RGBColor(0xF7, 0xF5, 0xF1)   # warm off-white, low glare
DARK   = RGBColor(0x1A, 0x1F, 0x24)   # story-mode background
TEAL   = RGBColor(0x0E, 0x7C, 0x86)
TEAL_L = RGBColor(0xDC, 0xED, 0xEF)
AMBER  = RGBColor(0xD9, 0x83, 0x24)
AMBER_L= RGBColor(0xFA, 0xEC, 0xD8)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RED_L  = RGBColor(0xF7, 0xDF, 0xDC)
GRAY   = RGBColor(0x6B, 0x76, 0x80)
GRAY_L = RGBColor(0xE4, 0xE2, 0xDE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xF7, 0xF5, 0xF1)

SANS = "Arial"
MONO = "Courier New"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.72)                       # side margin
CONTENT_W = W - 2 * M


# ---------------------------------------------------------------- core
class Deck:
    def __init__(self, tag):
        self.p = Presentation()
        self.p.slide_width, self.p.slide_height = W, H
        self.tag = tag
        self.n = 0

    def _blank(self, bg=GROUND):
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        r.fill.solid(); r.fill.fore_color.rgb = bg
        r.line.fill.background(); r.shadow.inherit = False
        return s

    def slide(self, bg=GROUND, footer=True, count=True):
        s = self._blank(bg)
        if count:
            self.n += 1
        if footer:
            dark = bg == DARK
            txt(s, self.tag, M, H - Inches(0.5), Inches(6), Inches(0.3), 11,
                GRAY if not dark else RGBColor(0x7A, 0x84, 0x8C))
            txt(s, str(self.n) if count else "", W - M - Inches(1), H - Inches(0.5),
                Inches(1), Inches(0.3), 11,
                GRAY if not dark else RGBColor(0x7A, 0x84, 0x8C), align=PP_ALIGN.RIGHT)
        return s

    def notes(self, s, text):
        s.notes_slide.notes_text_frame.text = text.strip()

    def save(self, path):
        self.p.save(path)
        return path


# ---------------------------------------------------------------- text
def txt(slide, text, x, y, w, h, size, color=INK, bold=False, font=SANS,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if line:
            para.line_spacing = line
        r = para.add_run(); r.text = ln
        f = r.font
        f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
        f.color.rgb = color
    return tb


def rich(slide, x, y, w, h, parts, size=20, align=PP_ALIGN.LEFT, line=1.25,
         space=6, anchor=MSO_ANCHOR.TOP):
    """parts: list of paragraphs; each is a list of (text, {opts}) runs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(parts):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        para.space_after = Pt(space)
        for text, o in para_runs:
            r = para.add_run(); r.text = text
            f = r.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", False)
            f.italic = o.get("italic", False)
            f.name = o.get("font", SANS)
            f.color.rgb = o.get("color", INK)
    return tb


# ---------------------------------------------------------------- furniture
def heading(slide, text, sub=None, dark=False, rule=TEAL, size=34):
    """Title with a thick accent rule beneath — the deck's signature."""
    c = CREAM if dark else INK
    txt(slide, text, M, Inches(0.52), CONTENT_W, Inches(0.85), size, c, bold=True)
    y = Inches(1.30) if not sub else Inches(1.30)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, y, Inches(1.5), Pt(4.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = rule
    ln.line.fill.background(); ln.shadow.inherit = False
    if sub:
        txt(slide, sub, M, Inches(1.48), CONTENT_W, Inches(0.4), 17,
            GRAY if not dark else RGBColor(0x9A, 0xA4, 0xAC), italic=True)
    return Inches(2.05) if sub else Inches(1.78)


def eyebrow(slide, text, color=TEAL, y=None):
    txt(slide, text.upper(), M, y or Inches(0.30), CONTENT_W, Inches(0.3), 12,
        color, bold=True)


def box(slide, x, y, w, h, text, fill=WHITE, edge=GRAY, tcolor=INK, size=15,
        bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE, edge_w=1.25, font=SANS,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, dash=None):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge; sh.line.width = Pt(edge_w)
    if dash:
        sh.line.dash_style = dash
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = 0.10
        except Exception:
            pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, ln in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run(); r.text = ln
        f = r.font
        f.size, f.bold, f.name = Pt(size), bold, font
        f.color.rgb = tcolor
    return sh


def arrow(slide, x, y, w, h=Inches(0.22), color=TEAL, label=None, lsize=12,
          up=True, label_y=None):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    if label:
        ly = label_y if label_y is not None else (
            y - Inches(0.30) if up else y + h + Inches(0.04))
        txt(slide, label, x - Inches(0.62), ly, w + Inches(1.24), Inches(0.28),
            lsize, color, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    return a


def down_arrow(slide, x, y, h, w=Inches(0.22), color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def statement(deck, text, sub=None, color=CREAM, bg=DARK, accent=AMBER,
              size=40, eyebrow_text=None):
    """Full-slide claim. Dark by default — the deck's punctuation."""
    s = deck.slide(bg=bg)
    if eyebrow_text:
        eyebrow(s, eyebrow_text, accent)
    txt(s, text, M, Inches(2.15), CONTENT_W, Inches(2.6), size, color, bold=True,
        line=1.18)
    if sub:
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(5.05), Inches(1.5), Pt(4))
        ln.fill.solid(); ln.fill.fore_color.rgb = accent
        ln.line.fill.background(); ln.shadow.inherit = False
        txt(s, sub, M, Inches(5.35), CONTENT_W - Inches(1), Inches(1.2), 19,
            RGBColor(0xB8, 0xC0, 0xC6), line=1.3)
    return s


def poll(deck, number, question, options, minute, correct=None, note=None,
         reveal=False):
    """Poll / ConcepTest slide. options: list of (letter, text)."""
    s = deck.slide()
    eyebrow(s, f"POLL {number}   ·   minute {minute}" + ("   ·   ANSWER" if reveal else ""),
            AMBER if not reveal else TEAL)
    # long stems need a smaller face and more room, or they clip
    L = len(question)
    qsize, y = (26, Inches(2.25)) if L <= 95 else (
        (23, Inches(2.50)) if L <= 165 else (20.5, Inches(2.72)))
    txt(s, question, M, Inches(0.72), CONTENT_W, y - Inches(0.80), qsize, INK,
        bold=True, line=1.2)
    n = len(options)
    hh = Inches(0.72) if n <= 4 else Inches(0.60)
    gap = Inches(0.16) if n <= 4 else Inches(0.11)
    for letter, text in options:
        hit = reveal and correct and letter == correct
        fill = TEAL_L if hit else WHITE
        edge = TEAL if hit else GRAY_L
        b = box(s, M, y, CONTENT_W, hh, "", fill=fill, edge=edge,
                edge_w=2.5 if hit else 1.25)
        txt(s, letter, M + Inches(0.22), y + Inches(0.13), Inches(0.5), hh,
            20, TEAL if hit else GRAY, bold=True)
        txt(s, text, M + Inches(0.78), y + Inches(0.14),
            CONTENT_W - Inches(1.1), hh, 18 if n <= 4 else 16,
            INK, bold=hit, anchor=MSO_ANCHOR.TOP)
        y += hh + gap
    if note:
        txt(s, note, M, Inches(6.06), CONTENT_W, Inches(0.78), 15,
            TEAL if reveal else GRAY, italic=True, line=1.25)
    return s


def section(deck, kicker, title, items=None, minute=None):
    s = deck.slide(bg=DARK)
    eyebrow(s, kicker, TEAL)
    txt(s, title, M, Inches(1.9), CONTENT_W - Inches(2), Inches(1.7), 44, CREAM,
        bold=True, line=1.15)
    if minute:
        txt(s, minute, W - M - Inches(3), Inches(0.28), Inches(3), Inches(0.4),
            13, GRAY, bold=True, font=MONO, align=PP_ALIGN.RIGHT)
    if items:
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(4.15), Inches(1.5), Pt(4))
        ln.fill.solid(); ln.fill.fore_color.rgb = TEAL
        ln.line.fill.background(); ln.shadow.inherit = False
        txt(s, "\n".join(items), M, Inches(4.55), CONTENT_W, Inches(2), 19,
            RGBColor(0xB8, 0xC0, 0xC6), line=1.5)
    return s


def table(slide, x, y, w, rows, col_w, head_fill=INK, size=15, row_h=Inches(0.44),
          head_size=None, mono_cols=(), align=None):
    """rows[0] is the header. col_w: fractions summing to 1."""
    nr, nc = len(rows), len(rows[0])
    t = slide.shapes.add_table(nr, nc, x, y, w, row_h * nr).table
    for j, frac in enumerate(col_w):
        t.columns[j].width = Emu(int(w * frac))
    for i, row in enumerate(rows):
        t.rows[i].height = row_h
        for j, cell_text in enumerate(row):
            c = t.cell(i, j)
            c.margin_left = c.margin_right = Inches(0.10)
            c.margin_top = c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = head_fill if i == 0 else (
                WHITE if i % 2 else GROUND)
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (align[j] if align else
                           (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
            r = p.add_run(); r.text = str(cell_text)
            f = r.font
            f.size = Pt(head_size or size) if i == 0 else Pt(size)
            f.bold = (i == 0) or (j == 0 and nc > 2)
            f.name = MONO if (j in mono_cols and i > 0) else SANS
            f.color.rgb = CREAM if i == 0 else INK
    return t


def chain(slide, y0, upto=7, highlight=(), values=None, dim_after=True,
          labels=None, size=13):
    """The measurement chain — the spine diagram of the whole course.

    7 stages, serpentine: 4 on the top row, 3 on the bottom.
    upto: how many stages to draw (progressive reveal).
    highlight: stage indices (1-based) drawn in amber.
    values: dict stage_index -> label on the arrow leaving that stage.
    """
    stages = labels or [
        "PHYSICAL\nQUANTITY",
        "MECHANICAL\nCOUPLING",
        "TRANSDUCTION",
        "ANALOG\nCONDITIONING",
        "SAMPLING &\nQUANTISATION",
        "CODES → SI\nUNITS",
        "TIMESTAMP,\nLOG, DECIDE",
    ]
    bw, bh, gap = Inches(2.42), Inches(1.02), Inches(0.52)
    row1, row2 = stages[:4], stages[4:]
    positions = []
    for i, _ in enumerate(row1):
        positions.append((M + i * (bw + gap), y0))
    y1 = y0 + bh + Inches(0.92)
    # second row flows right-to-left, so the diagram reads as a snake
    for i, _ in enumerate(row2):
        positions.append((M + (3 - i) * (bw + gap), y1))

    for idx, (label, (x, y)) in enumerate(zip(stages, positions), start=1):
        if idx > upto:
            if not dim_after:
                continue
            box(slide, x, y, bw, bh, label, fill=GROUND, edge=GRAY_L,
                tcolor=GRAY_L, size=size, bold=True,
                dash=MSO_LINE_DASH_STYLE.DASH)
            continue
        hot = idx in highlight
        box(slide, x, y, bw, bh, label,
            fill=AMBER_L if hot else WHITE,
            edge=AMBER if hot else TEAL,
            tcolor=INK, size=size, bold=True, edge_w=2.5 if hot else 1.5)
        txt(slide, str(idx), x + Inches(0.08), y + Inches(0.03), Inches(0.3),
            Inches(0.25), 10, AMBER if hot else TEAL, bold=True, font=MONO)

    # arrows between boxes: right-pointing on the top row, left-pointing below
    for idx in range(1, min(upto, 7)):
        if idx == 4:      # wrap down from box 4 to box 5, same column
            ax = M + 3 * (bw + gap) + bw / 2
            down_arrow(slide, ax - Inches(0.11), y0 + bh + Inches(0.14),
                       Inches(0.62), color=TEAL)
            continue
        lbl = (values or {}).get(idx)
        if idx < 4:                                   # top row, →
            x = M + (idx - 1) * (bw + gap) + bw
            arrow(slide, x + Inches(0.06), y0 + bh / 2 - Inches(0.11),
                  gap - Inches(0.12), color=TEAL, label=lbl,
                  label_y=y0 - Inches(0.34))
        else:                                         # bottom row, ←
            col = 3 - (idx - 4)                       # box idx sits in this column
            x = M + (col - 1) * (bw + gap) + bw
            a = slide.shapes.add_shape(
                MSO_SHAPE.LEFT_ARROW, x + Inches(0.06),
                y1 + bh / 2 - Inches(0.11), gap - Inches(0.12), Inches(0.22))
            a.fill.solid(); a.fill.fore_color.rgb = TEAL
            a.line.fill.background(); a.shadow.inherit = False
            if lbl:
                txt(slide, lbl, x - Inches(0.56), y1 - Inches(0.34),
                    gap + Inches(1.12), Inches(0.28), 12, TEAL, bold=True,
                    font=MONO, align=PP_ALIGN.CENTER)
    return y1 + bh


def curve(slide, pts, color=TEAL, width=1.75, dash=None):
    """Polyline through pts [(x,y), ...] as EMU lengths — used for waveforms."""
    b = slide.shapes.build_freeform(pts[0][0], pts[0][1])
    b.add_line_segments(pts[1:], close=False)
    sh = b.convert_to_shape()
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(width)
    if dash:
        sh.line.dash_style = dash
    sh.shadow.inherit = False
    return sh


def dot(slide, cx, cy, r=Inches(0.055), color=RED):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r)
    d.fill.solid(); d.fill.fore_color.rgb = color
    d.line.color.rgb = WHITE; d.line.width = Pt(0.75)
    d.shadow.inherit = False
    return d


def axis(slide, x, y, w, color=GRAY_L, width=1.0):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(width))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def sine(x0, y_mid, w, amp, cycles, n=400, phase=0.0):
    """Sample points for a sine wave, as EMU tuples."""
    import math
    return [(Emu(int(x0 + w * i / n)),
             Emu(int(y_mid - amp * math.sin(2 * math.pi * cycles * i / n + phase))))
            for i in range(n + 1)]

# -*- coding: utf-8 -*-
"""Produce the Armenian decks from the built English ones.

    cd lectures/tools && .venv/bin/python translate_deck.py

Reads   ../lecture-0{1,2}/output/L{1,2}-*.pptx
Writes  ../lecture-0{1,2}/output/L{1,2}-*-HY.pptx

Nothing is translated in place: the English decks are untouched, so both languages
stay in step whenever build_l1.py / build_l2.py are re-run.

Fonts. Arial and Courier New carry no Armenian glyphs, so any run that ends up
containing Armenian is switched to DejaVu Sans (or DejaVu Sans Mono if it was
monospaced). Runs that stay Latin — register values, part numbers, units — keep
their original font, which preserves their digit metrics.

Note on Noto: "Noto Sans Armenian" was rejected deliberately. It has no Latin
digits, so LibreOffice substitutes them glyph by glyph and every number renders
with gaps ("0 . 0 6 1"). In a course made of numbers that is unusable.
"""
import glob
import os
import re
import sys

from pptx import Presentation
from pptx.util import Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from i18n import hy_part1, hy_part2, hy_part3, hy_part4, hy_part5  # noqa: E402

HY = {}
for mod in (hy_part1, hy_part2, hy_part3, hy_part4, hy_part5):
    HY.update(mod.HY)

ARM = re.compile(r"[԰-֏ﬓ-ﬗ]")
# strings that are pure number / unit / symbol need no translation and no font change
SKIP = re.compile(r"^[\s\d.,;:%×√±°µ/()|\-–—+=<>x\[\]{}A-Fa-f0-9€$≥≤≈∝²³₀]*$")

SANS_HY, MONO_HY = "DejaVu Sans", "DejaVu Sans Mono"

missing, translated = [], 0


def font_for(original_name, text):
    """Pick a font that can actually render `text`."""
    if not ARM.search(text or ""):
        return original_name                      # still Latin: leave it alone
    if original_name and "Courier" in original_name:
        return MONO_HY
    return SANS_HY


def do_paragraph(para):
    """Translate a paragraph, collapsing its runs into one.

    Armenian word order rarely matches English, so translating run-by-run would
    scramble the sentence. Instead the whole paragraph is looked up as one string;
    the first run keeps its formatting and inherits bold if ANY run was bold.
    """
    global translated
    runs = para.runs
    if not runs:
        return
    whole = "".join(r.text for r in runs)
    key = whole.strip()
    if not key:
        return

    if key in HY:
        new = HY[key]
        translated += 1
    elif SKIP.match(key):
        new = None                                # numeric: nothing to do
    else:
        missing.append(key)
        new = None

    keep = runs[0]
    if new is not None:
        any_bold = any(r.font.bold for r in runs)
        # preserve leading/trailing whitespace of the original
        lead = whole[: len(whole) - len(whole.lstrip())]
        trail = whole[len(whole.rstrip()):]
        keep.text = lead + new + trail
        for r in runs[1:]:
            r.text = ""
        if any_bold:
            keep.font.bold = True

    # Armenian runs ~10-25 % longer than English. Where a translated string grew,
    # ease the point size down a little so it still fits a fixed-width box.
    if new is not None and keep.font.size and len(key) > 3:
        grow = len(new) / max(len(key), 1)
        if grow > 1.12:
            factor = max(0.80, min(1.0, (1.0 / grow) ** 0.55))
            keep.font.size = Pt(round(keep.font.size.pt * factor, 1))

    keep.font.name = font_for(keep.font.name, keep.text)
    for r in runs[1:]:
        if r.text:
            r.font.name = font_for(r.font.name, r.text)


def walk(shape):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            do_paragraph(para)
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    do_paragraph(para)
    if shape.shape_type == 6:                     # group
        for sub in shape.shapes:
            walk(sub)


def translate(path):
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            walk(shape)
        # speaker notes: instructor-facing, kept in English on purpose (see README)
    out = path.replace(".pptx", "-HY.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    for pattern in ("../lecture-01/output/L1-*.pptx", "../lecture-02/output/L2-*.pptx"):
        for f in sorted(glob.glob(pattern)):
            if f.endswith("-HY.pptx"):
                continue
            print("→", translate(f))

    uniq = sorted(set(missing))
    print(f"\ntranslated paragraphs : {translated}")
    print(f"untranslated strings   : {len(uniq)}")
    if uniq:
        with open("i18n/_untranslated.txt", "w") as fh:
            fh.write("\n".join(uniq))
        for s in uniq[:40]:
            print("   ·", s[:100])
        if len(uniq) > 40:
            print(f"   … {len(uniq) - 40} more in i18n/_untranslated.txt")
